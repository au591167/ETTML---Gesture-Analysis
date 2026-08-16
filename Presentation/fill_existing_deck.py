from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PRESENTATION_DIR = ROOT / "Presentation"
ASSETS = PRESENTATION_DIR / "assets"
REPORT_IMAGES = ROOT / "Report" / "src" / "img"
SOURCE = PRESENTATION_DIR / "ETTML-Præsentation.pptx"
OUTPUT = PRESENTATION_DIR / "ETTML-Præsentation-final.pptx"

NAVY = RGBColor(9, 31, 63)
TEXT = RGBColor(31, 48, 73)
BLUE = RGBColor(47, 103, 224)
GREEN = RGBColor(22, 163, 74)
ORANGE = RGBColor(224, 119, 0)
PURPLE = RGBColor(124, 58, 237)
RED = RGBColor(220, 38, 38)
LIGHT = RGBColor(246, 249, 252)
PALE_BLUE = RGBColor(226, 238, 252)
PALE_GREEN = RGBColor(228, 248, 236)
PALE_ORANGE = RGBColor(255, 244, 216)
PALE_PURPLE = RGBColor(239, 233, 254)
PALE_RED = RGBColor(254, 232, 232)
MUTED = RGBColor(92, 108, 130)
WHITE = RGBColor(255, 255, 255)


def remove_all_shapes(slide):
    for shape in list(slide.shapes):
        element = shape._element
        element.getparent().remove(element)


def set_background(slide, color=LIGHT):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, text, x, y, w, h, size=22, color=TEXT, bold=False,
             align=PP_ALIGN.LEFT, font="Aptos", valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def add_title(slide, title, subtitle=None):
    add_text(slide, title, 0.72, 0.32, 11.9, 0.56, 29, NAVY, True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(1.02), Inches(1.05), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()
    if subtitle:
        add_text(slide, subtitle, 1.95, 0.87, 10.6, 0.28, 12, MUTED)


def add_card(slide, x, y, w, h, fill_color, border_color, title, body="",
             title_size=21, body_size=15, align=PP_ALIGN.CENTER):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    card.line.color.rgb = border_color
    card.line.width = Pt(1.7)
    tf = card.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.11)
    tf.margin_bottom = Inches(0.08)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = align
    p.font.name = "Aptos Display"
    p.font.size = Pt(title_size)
    p.font.bold = True
    p.font.color.rgb = NAVY
    if body:
        p2 = tf.add_paragraph()
        p2.text = body
        p2.alignment = align
        p2.space_before = Pt(7)
        p2.font.name = "Aptos"
        p2.font.size = Pt(body_size)
        p2.font.color.rgb = TEXT
    return card


def add_picture_contain(slide, path, x, y, w, h):
    with Image.open(path) as img:
        iw, ih = img.size
    region_ratio = w / h
    image_ratio = iw / ih
    if image_ratio > region_ratio:
        pw = w
        ph = w / image_ratio
        px = x
        py = y + (h - ph) / 2
    else:
        ph = h
        pw = h * image_ratio
        px = x + (w - pw) / 2
        py = y
    return slide.shapes.add_picture(str(path), Inches(px), Inches(py), Inches(pw), Inches(ph))


def use_full_asset(slide, path):
    remove_all_shapes(slide)
    set_background(slide)
    slide.shapes.add_picture(str(path), 0, 0, width=Inches(13.333333), height=Inches(7.5))


def add_arrow(slide, x, y, w=0.75, h=0.42, color=MUTED):
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()
    return arrow


def add_code_block(slide, code, x, y, w, h, size=13, caption="kodeudsnit"):
    panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(15, 23, 42)
    panel.line.color.rgb = RGBColor(51, 65, 85)
    panel.line.width = Pt(1.2)
    panel.text_frame.clear()

    for i, color in enumerate((RED, ORANGE, GREEN)):
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x + 0.22 + i * 0.18), Inches(y + 0.19),
            Inches(0.10), Inches(0.10)
        )
        dot.fill.solid(); dot.fill.fore_color.rgb = color
        dot.line.fill.background()
    add_text(slide, caption, x + 0.82, y + 0.12, w - 1.05, 0.24,
             10, RGBColor(148, 163, 184), False, PP_ALIGN.LEFT, "Consolas")
    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x + 0.16), Inches(y + 0.46),
        Inches(w - 0.32), Inches(0.015)
    )
    divider.fill.solid(); divider.fill.fore_color.rgb = RGBColor(51, 65, 85)
    divider.line.fill.background()

    code_box = add_text(
        slide, code, x + 0.28, y + 0.62, w - 0.56, h - 0.76,
        size, RGBColor(226, 232, 240), False, PP_ALIGN.LEFT, "Consolas",
        MSO_ANCHOR.TOP,
    )
    code_box.text_frame.word_wrap = False
    code_box.text_frame.margin_left = 0
    code_box.text_frame.margin_right = 0
    code_box.text_frame.margin_top = 0
    code_box.text_frame.margin_bottom = 0
    code_box.text_frame.paragraphs[0].line_spacing = 1.02
    return panel


def move_last_slide(prs, target_index):
    """Move the slide most recently appended to a zero-based target index."""
    slide_id = prs.slides._sldIdLst[-1]
    prs.slides._sldIdLst.remove(slide_id)
    prs.slides._sldIdLst.insert(target_index, slide_id)


prs = Presentation(str(SOURCE))

# Slide 1 — title
s = prs.slides[0]
remove_all_shapes(s)
set_background(s, WHITE)
photo = REPORT_IMAGES / "Circuit.jpg"
add_picture_contain(s, photo, 6.85, 0, 6.48, 7.5)
overlay = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(7.25), Inches(7.5))
overlay.fill.solid()
overlay.fill.fore_color.rgb = NAVY
overlay.line.fill.background()
add_text(s, "Tiny Machine\nLearning", 0.8, 1.05, 5.6, 1.55, 38, WHITE, True, font="Aptos Display")
add_text(s, "Gesture Recognition on the Edge", 0.83, 2.82, 5.7, 0.5, 21, RGBColor(147, 197, 253), True)
add_text(s, "Particle Photon 2 + ADXL343", 0.83, 3.48, 5.7, 0.38, 16, WHITE)
add_text(s, "Eksamenspræsentation", 0.83, 4.58, 5.7, 0.35, 15, RGBColor(203, 213, 225))
add_text(s, "Erik Kjær Klint  ·  201704536", 0.83, 5.02, 5.8, 0.4, 17, WHITE, True)

# Slide 2 — agenda
s = prs.slides[1]
remove_all_shapes(s)
set_background(s)
add_title(s, "Agenda", "En kronologisk rejse fra fysisk bevægelse til lokal handling")
agenda = [
    ("01", "Projekt + demo", PALE_BLUE, BLUE),
    ("02", "Hardware + data", PALE_GREEN, GREEN),
    ("03", "Features", PALE_ORANGE, ORANGE),
    ("04", "Model + deployment", PALE_PURPLE, PURPLE),
    ("05", "Resultater", PALE_RED, RED),
]
for i, (num, label, fill, border) in enumerate(agenda):
    x = 0.67 + i * 2.52
    add_card(s, x, 2.25, 2.18, 2.35, fill, border, num, label, 30, 16)
    if i < len(agenda) - 1:
        add_arrow(s, x + 2.2, 3.18, 0.32, 0.25)
add_text(s, "PROBLEM  →  DATA  →  FEATURES  →  MODEL  →  RESULTAT", 1.1, 5.35, 11.1, 0.55, 20, NAVY, True, PP_ALIGN.CENTER)

# Slides built from the prepared full-slide visual assets.
full_assets = {
    2: "gesture-command-map.png",       # slide 3
    4: "hardware-wiring.png",           # slide 5
    5: "system-dataflow.png",           # slide 6
    6: "data-acquisition-overview.png", # slide 7
    7: "quality-control-comparison.png",# slide 8
    8: "feature-pipeline.png",           # slide 9
    9: "mlp-architecture.png",           # slide 10
    10: "dataset-split-cv.png",          # slide 11
    11: "deployment-decision-flow.png",  # slide 12
}
for slide_index, filename in full_assets.items():
    use_full_asset(prs.slides[slide_index], ASSETS / filename)

# Replace the densest evaluation heading with plain language while retaining
# the technical term as a small translation underneath.
s = prs.slides[10]
cv_cover = s.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(3.82), Inches(6.15), Inches(0.64)
)
cv_cover.fill.solid()
cv_cover.fill.fore_color.rgb = WHITE
cv_cover.line.fill.background()
add_text(s, "5 testrunder med skiftende testgruppe", 0.88, 3.88, 5.75, 0.32, 18, NAVY, True)
add_text(s, "Teknisk: stratificeret 5-fold cross-validation", 0.88, 4.18, 5.75, 0.20, 10, MUTED)

# Slide 4 — demo
s = prs.slides[3]
remove_all_shapes(s)
set_background(s)
add_title(s, "Demo", "Fra fysisk gestus til lokal kommando")
photo_card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.35), Inches(6.15), Inches(5.35))
photo_card.fill.solid(); photo_card.fill.fore_color.rgb = WHITE
photo_card.line.color.rgb = RGBColor(203, 213, 225); photo_card.line.width = Pt(1.5)
add_picture_contain(s, photo, 0.86, 1.53, 5.83, 4.55)
add_text(s, "Photon 2 klassificerer lokalt — ingen cloud-inference", 1.0, 6.08, 5.55, 0.34, 14, NAVY, True, PP_ALIGN.CENTER)
steps = [
    ("1", "Udfør gestus", "tap1 · tap2 · tap3 · shake", PALE_BLUE, BLUE),
    ("2", "Lokal klassifikation", "features → scaler → MLP", PALE_PURPLE, PURPLE),
    ("3", "Observerbart output", "EVENT + RGB-feedback", PALE_GREEN, GREEN),
]
for i, (num, title, body, fill, border) in enumerate(steps):
    y = 1.52 + i * 1.63
    add_card(s, 7.25, y, 5.25, 1.28, fill, border, f"{num}   {title}", body, 20, 14, PP_ALIGN.LEFT)
    if i < 2:
        add_text(s, "↓", 9.55, y + 1.24, 0.6, 0.35, 20, MUTED, True, PP_ALIGN.CENTER)
add_text(s, "Demoen viser produktet først — derefter følger forklaringen dataene fra sensoren.", 7.32, 6.35, 5.0, 0.35, 13, MUTED, False, PP_ALIGN.CENTER)

# Slide 13 — measured results
s = prs.slides[12]
remove_all_shapes(s)
set_background(s)
add_title(s, "Resultater", "Rapportens dokumenterede 25-optagelses-snapshot")
metrics = [
    ("80 %", "holdout accuracy\n4 af 5 korrekte", PALE_BLUE, BLUE),
    ("76 %", "5-fold CV\n19 af 25 korrekte", PALE_GREEN, GREEN),
    ("345 µs", "gennemsnitlig\nmodelinferens", PALE_PURPLE, PURPLE),
    ("0", "registrerede\nsensor-read-fejl", PALE_ORANGE, ORANGE),
]
for i, (value, label, fill, border) in enumerate(metrics):
    x = 0.67 + (i % 2) * 2.5
    y = 1.45 + (i // 2) * 2.1
    add_card(s, x, y, 2.22, 1.72, fill, border, value, label, 27, 14)
add_text(s, "Macro precision 70,0 %  ·  recall 80,0 %  ·  F1 73,3 %", 0.72, 5.78, 4.63, 0.48, 13, TEXT, True, PP_ALIGN.CENTER)
add_text(s, "Én fejl i fem testeksempler ændrer accuracy med 20 procentpoint.", 0.72, 6.28, 4.63, 0.38, 12, RED, True, PP_ALIGN.CENTER)
cm_card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.58), Inches(1.28), Inches(7.08), Inches(5.48))
cm_card.fill.solid(); cm_card.fill.fore_color.rgb = WHITE
cm_card.line.color.rgb = RGBColor(203, 213, 225); cm_card.line.width = Pt(1.3)
add_picture_contain(s, REPORT_IMAGES / "confusion_matrix.png", 5.73, 1.42, 6.78, 5.18)

# Slide 14 — current data status and next work
layout = prs.slide_layouts[6]
s = prs.slides.add_slide(layout)
set_background(s)
add_title(s, "Aktuel datastatus", "Resultater: videre arbejde holdes tydeligt adskilt fra rapportens metrics")
add_card(s, 0.72, 1.48, 4.65, 3.15, PALE_BLUE, BLUE,
         "Rapportens snapshot · 25",
         "5 optagelser per klasse\nBalanceret og deployeret\nGrundlag for de rapporterede metrics", 25, 17)
add_arrow(s, 5.55, 2.75, 1.05, 0.55, PURPLE)
add_card(s, 6.82, 1.48, 5.78, 3.15, PALE_PURPLE, PURPLE,
         "Aktuel kandidatpulje · 49",
         "idle 5 · tap1 5 · tap2 14\ntap3 20 · shake 5\nKontraktkompatibel, men ikke balanceret", 25, 17)
add_text(s, "Næste kontrollerede iteration", 0.75, 5.05, 12.0, 0.42, 20, NAVY, True, PP_ALIGN.CENTER)
next_steps = [
    ("1", "Indsaml manglende klasser"),
    ("2", "Kurater og balancér"),
    ("3", "Retræn modellen"),
    ("4", "Ny blind evaluering"),
]
for i, (num, label) in enumerate(next_steps):
    x = 0.74 + i * 3.14
    add_card(s, x, 5.62, 2.73, 0.93, WHITE, [BLUE, GREEN, ORANGE, RED][i], f"{num}  {label}", "", 16, 12)
    if i < 3:
        add_arrow(s, x + 2.76, 5.89, 0.34, 0.28)
add_text(s, "Konklusion: fungerende end-to-end-prototype · begrænset evidens for generalisering", 0.75, 6.78, 11.85, 0.35, 15, NAVY, True, PP_ALIGN.CENTER)

# New slide placed directly after the hardware slide. It tells one focused
# story: an initial axis assumption was corrected through measured data.
s = prs.slides.add_slide(layout)
set_background(s)
add_title(
    s,
    "Fra antagelse til målt orientering",
    "Start: Y op/ned  →  faktisk: Y venstre/højre · X frem/tilbage · Z op/ned",
)

visual_card = s.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(1.24), Inches(8.30), Inches(5.62)
)
visual_card.fill.solid(); visual_card.fill.fore_color.rgb = WHITE
visual_card.line.color.rgb = RGBColor(203, 213, 225); visual_card.line.width = Pt(1.3)
add_picture_contain(s, ASSETS / "adxl343-axis-orientation.png", 0.70, 1.38, 8.00, 4.50)
add_text(s, "Datastrømmen var hele tiden X, Y, Z — jeg rettede den fysiske fortolkning.",
         0.82, 6.18, 7.76, 0.34, 13, NAVY, True, PP_ALIGN.CENTER)

add_text(s, "Hvad lærte jeg?", 9.18, 1.33, 3.40, 0.40, 20, NAVY, True, PP_ALIGN.CENTER)
learning = [
    ("1 · ANTAGELSE", "Y var op/ned", PALE_ORANGE, ORANGE),
    ("2 · MÅLING", "Z-RMS ≈ 2 × Y-RMS", PALE_BLUE, BLUE),
    ("3 · DESIGN", "X + Y + Z + magnitude", PALE_GREEN, GREEN),
]
for i, (title, body, fill, border) in enumerate(learning):
    y = 1.90 + i * 1.57
    add_card(s, 9.05, y, 3.60, 1.18, fill, border, title, body, 16, 14)
    if i < 2:
        add_text(s, "↓", 10.50, y + 1.16, 0.70, 0.38, 22, MUTED, True, PP_ALIGN.CENTER)

# python-pptx appends new slides; move this one from the end to position 6.
move_last_slide(prs, 5)

# Code focus 1: feature extraction. Placed directly after the feature diagram.
s = prs.slides.add_slide(layout)
set_background(s)
add_title(s, "Kodefokus · Fra vindue til 28 features",
          "Product/ml/train.py · samme featureorden skal bruges i Python og C++")
feature_code = (
    "# Vælg de fire signal-kanaler\n"
    'channels = ["ax", "ay", "az"]\n'
    "if use_mag:\n"
    '    channels = channels + ["mag"]\n\n'
    "feats = []\n"
    "for c in channels:\n"
    "    x = window_df[c].to_numpy(dtype=np.float32)\n"
    "    # Fjern kanalens statiske offset\n"
    "    x = x - np.mean(x)\n"
    "    # Tilføj syv tal til featurevektoren\n"
    "    feats.extend(channel_features(x))\n\n"
    "# channel_features returnerer:\n"
    "[std, mn, mx, rng, energy,\n"
    " peaks, max_abs_diff]"
)
add_code_block(s, feature_code, 0.70, 1.40, 7.22, 4.75, 12,
               "train.py  ·  extract_features_for_window()")
add_text(s, "Faktisk princip · forkortet uden at ændre rækkefølgen", 0.86, 6.25, 6.90, 0.28,
         11, MUTED, False, PP_ALIGN.CENTER)
feature_steps = [
    ("INPUT", "1.600 samples\nX, Y, Z + magnitude", PALE_BLUE, BLUE),
    ("BEHANDLING", "mean removal\n7 statistikker", PALE_ORANGE, ORANGE),
    ("OUTPUT", "4 × 7\n= 28 modelinputs", PALE_GREEN, GREEN),
]
for i, (title, body, fill, border) in enumerate(feature_steps):
    y = 1.52 + i * 1.62
    add_card(s, 8.42, y, 4.15, 1.25, fill, border, title, body, 17, 14)
    if i < 2:
        add_text(s, "↓", 10.13, y + 1.22, 0.70, 0.35, 21, MUTED, True, PP_ALIGN.CENTER)
add_text(s, "KODEKONTRAKT · ændret featureorden = forkert input til de lærte vægte",
         8.35, 6.43, 4.30, 0.34, 11, RED, True, PP_ALIGN.CENTER)
move_last_slide(prs, 10)

# Code focus 2: split, scaling, and model fitting. Placed after evaluation.
s = prs.slides.add_slide(layout)
set_background(s)
add_title(s, "Kodefokus · Split, skalering og træning",
          "Product/ml/train.py · de centrale linjer der beskytter evalueringen")
training_code = (
    "# Del data: 80 % træning og 20 % test\n"
    "X_train, X_test, y_train, y_test = train_test_split(\n"
    "    X, y_enc,\n"
    "    test_size=0.2,\n"
    "    random_state=42,\n"
    "    stratify=y_enc,\n"
    ")\n\n"
    "# Lær skaleringen KUN fra træningsdata\n"
    "scaler = StandardScaler()\n"
    "X_train_s = scaler.fit_transform(X_train)\n"
    "X_test_s  = scaler.transform(X_test)\n\n"
    "# Træn modellen og test på ukendte data\n"
    "model = train_model(cfg, X_train_s, y_train)\n"
    "y_pred = model.predict(X_test_s)"
)
add_code_block(s, training_code, 0.70, 1.40, 7.28, 4.92, 11,
               "train.py  ·  main()")
add_text(s, "model.fit(...) inde i train_model() udfører træning og backpropagation",
         0.82, 6.38, 7.05, 0.32, 11, MUTED, True, PP_ALIGN.CENTER)
training_points = [
    ("STRATIFY", "én testoptagelse\nfra hver klasse", PALE_BLUE, BLUE),
    ("SEED 42", "samme split kan\ngentages", PALE_PURPLE, PURPLE),
    ("FIT KUN TRAIN", "testdata påvirker ikke\nscalerens μ og σ", PALE_GREEN, GREEN),
    ("20 → 5", "træn på 20\ntest på 5", PALE_ORANGE, ORANGE),
]
for i, (title, body, fill, border) in enumerate(training_points):
    x = 8.35 + (i % 2) * 2.18
    y = 1.56 + (i // 2) * 2.08
    add_card(s, x, y, 1.98, 1.64, fill, border, title, body, 14, 12)
add_text(s, "Ingen leakage: fit_transform(train) · kun transform(test)",
         8.34, 5.90, 4.24, 0.52, 12, RED, True, PP_ALIGN.CENTER)
move_last_slide(prs, 13)

# Code focus 3: deployed inference and decision logic. Placed after deployment.
s = prs.slides.add_slide(layout)
set_background(s)
add_title(s, "Kodefokus · Fra 28 features til en sikker EVENT",
          "Product/firmware/src/main.cpp + model_data.cpp · inference er mere end modelkaldet")
runtime_code = (
    "// Komprimér sensorvinduet til 28 features\n"
    "extractFeatures(features);\n"
    "// Forward-pass giver fem klassescores\n"
    "tinyml_model::model_infer(\n"
    "    features, 28, scores, 5);\n\n"
    "// Find klassen med den højeste score\n"
    "int bestIdx = 0;\n"
    "float bestScore = scores[0];\n"
    "for (size_t i = 1; i < 5; ++i) {\n"
    "    if (scores[i] > bestScore) {\n"
    "        bestScore = scores[i];\n"
    "        bestIdx = (int)i;\n"
    "    }\n"
    "}\n\n"
    "// Kræv sikker og stabil prediction\n"
    "const int eventClass = updateDecision(\n"
    "    bestIdx, bestScore, millis());\n"
    "if (eventClass >= 0) { /* EVENT + RGB */ }"
)
add_code_block(s, runtime_code, 0.70, 1.40, 7.12, 4.88, 11,
               "main.cpp  ·  runInference()")
add_text(s, "model_infer(): StandardScaler → vægtede summer → ReLU → softmax",
         0.78, 6.36, 6.98, 0.34, 11, MUTED, True, PP_ALIGN.CENTER)
runtime_steps = [
    ("28 features", PALE_ORANGE, ORANGE),
    ("scaler + MLP", PALE_PURPLE, PURPLE),
    ("5 softmax-scores", PALE_BLUE, BLUE),
    ("score ≥ 0,75", PALE_GREEN, GREEN),
    ("3 ens + 4 s debounce", PALE_RED, RED),
    ("EVENT + RGB", WHITE, NAVY),
]
for i, (label, fill, border) in enumerate(runtime_steps):
    y = 1.32 + i * 0.88
    add_card(s, 8.42, y, 4.12, 0.66, fill, border, label, "", 14, 11)
    if i < len(runtime_steps) - 1:
        add_text(s, "↓", 10.12, y + 0.62, 0.70, 0.27, 16, MUTED, True, PP_ALIGN.CENTER)
add_text(s, "Photon 2: forward-pass med faste vægte · ingen backpropagation",
         8.35, 6.66, 4.27, 0.30, 11, NAVY, True, PP_ALIGN.CENTER)
move_last_slide(prs, 15)

# A consistent glanceable cue line for the presenter. These are deliberately
# short enough to help recover the sequence without turning slides into notes.
guide_lines = {
    3: "Formål  →  fem klasser  →  kommandoer",
    4: "Gestus  →  lokal klassifikation  →  EVENT + RGB",
    5: "Hardware  →  I2C  →  rå X/Y/Z-data",
    6: "Antagelse: Y op/ned  →  faktisk: Y venstre/højre · X frem/tilbage · Z op/ned",
    7: "Følg dataflowet fra venstre mod højre",
    8: "Baseline og støj  →  400 Hz  →  1.600 samples  →  kvalitet",
    9: "Godkend eller afvis  →  bedre labels  →  risiko for for rene data",
    10: "Mean removal + magnitude  →  4 × 7  →  28 features",
    11: "Kode: 4 kanaler  →  mean removal  →  7 features  →  28 inputs",
    12: "Skalering  →  neuroner  →  ReLU  →  softmax",
    13: "20 train + 5 test  →  seed 42  →  gentag 5 gange  →  76 %",
    14: "Kode: stratify  →  seed 42  →  scaler fit kun på train  →  model.fit",
    15: "Offline træning  →  C++ forward-pass  →  stabil handling",
    16: "Kode: features  →  scores  →  threshold  →  stabilitet  →  EVENT",
    17: "Resultat  →  fortolkning  →  lille testset  →  begrænsning",
    18: "25 snapshot  →  49 kandidater  →  balancér  →  retræn",
}
for slide_number, guide in guide_lines.items():
    slide = prs.slides[slide_number - 1]
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.12), Inches(13.333333), Inches(0.38)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    add_text(
        slide, f"GUIDE  ·  {guide}", 0.25, 7.145, 12.83, 0.25,
        11, WHITE, True, PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE
    )

prs.save(str(OUTPUT))
print(OUTPUT)
